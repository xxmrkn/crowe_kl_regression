import os
import tqdm
import glob
from itertools import product

import numpy as np
import pandas as pd

import cv2
import tqdm
import imageio

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.dml import MSO_THEME_COLOR


def main():
    path = 'Y:/masuda/project/crowe_kl_regression_update/report/1111_results/csvfile/VisionTransformer_Base16/results_csv2.csv'
    pptx_path = '/outliers.pptx' 
    _ROOT_TGT_PPT = 'Y:/masuda/project/crowe_kl_regression_update/report/1111_results'
    _ROOT_IMG = 'Y:/masuda/project/vit_kl_crowe/20220511_DRR_with_Crowe_KL/DRR_AP'

    def read_list_file(fpath):
        datalist = []
        if fpath.endswith('.txt'): #textfile
            with open(fpath, 'r' ,encoding="utf-8_sig", errors='ignore') as fr:
                datalist = fr.readlines()
            datalist = [dat.replace('\n', '') for dat in datalist]
        elif fpath.endswith('.csv'): #csv file
            df = pd.read_csv(fpath)
            print('Dataframe: ', df)
            print(datalist)
        return datalist

    def resize_img(img):
        scale_percent = 30 # percent of original size
        width = int(img.shape[1] * scale_percent / 100)
        height = int(img.shape[0] * scale_percent / 100)
        dim = (width, height)
        return cv2.resize(img, dim, interpolation = cv2.INTER_AREA)

    def chunks(l, n):
        n = max(1, n)
        return (l[i:i+n] for i in range(0, len(l), n))

    os.makedirs(_ROOT_TGT_PPT,exist_ok=True) 
    _ROWS = 5
    _COLS = 10

    _tops = [Cm(i) for i in np.arange(0,_ROWS*5,4.3)]
    _tops_txt = [Cm(i) for i in np.arange(2.8,_ROWS*5,4.3)]
    _ht = Cm(3)
    _lefts_vol = [Cm(i) for i in np.arange(0,_COLS*3.1,3.1)]
    _lefts_txt = [Cm(i) for i in np.arange(0,_COLS*3.1,3.1)]

    _MAX_SLIDES = 10
    _MAX_CASES  = 3000
    _EXC_LIST_F = ''
    _EXC_LIST = read_list_file(_EXC_LIST_F)
    _IMG_EXT   = '_AP.png'

    _LOG_FILE = os.path.join(_ROOT_TGT_PPT, 'log.txt')
    _TMP_IMG = 'tmp_strata.png'
    with open(_LOG_FILE, 'w') as f:
        f.writelines('Unprocessed list\n')

    _EXTS = None
    _CASE_LIST_F = path
    print(f'case_list {_CASE_LIST_F}')
    _CASE_LIST = pd.read_csv(_CASE_LIST_F, header=0, index_col=0)
    print(f'_CASE_LIST{_CASE_LIST}')

    _CASE_LIST = _CASE_LIST[abs(_CASE_LIST['Diff'])>=2]
    print(_CASE_LIST)

    _CASE_LIST = _CASE_LIST.replace({'True':{0:'Crowe=1,KL=1',
                                             1:'Crowe=1,KL=2',
                                             2:'Crowe=1,KL=3',
                                             3:'Crowe=1,KL=4',
                                             4:'Crowe=2,KL=4',
                                             5:'Crowe=3,KL=4',
                                             6:'Crowe=4,KL=4'}})
    _CASE_LIST = _CASE_LIST.replace({'Pred':{0:'Crowe=1,KL=1',
                                             1:'Crowe=1,KL=2',
                                             2:'Crowe=1,KL=3',
                                             3:'Crowe=1,KL=4',
                                             4:'Crowe=2,KL=4',
                                             5:'Crowe=3,KL=4',
                                             6:'Crowe=4,KL=4'}})

    imgs = _CASE_LIST.index.tolist()

    _IMG_CHUNKS = chunks(imgs,_MAX_CASES)
    for h,_IMG_CHUNK in enumerate(_IMG_CHUNKS):
        prs = Presentation()
        prs.slide_width    = 11887200
        prs.slide_height   = 7786550
        title_slide_layout = prs.slide_layouts[5]
        blank_slide_layout = prs.slide_layouts[6]
        print('Image chunk #', h)
        _SLIDE_CHUNKS = chunks(_IMG_CHUNK,_ROWS*_COLS)
        pbar = tqdm.tqdm(_SLIDE_CHUNKS)
        for i,_IMAGE_LIST_CHUNK in enumerate(pbar):
            print('\tSlide chunk #', i)
            slide = prs.slides.add_slide(blank_slide_layout)
            _SUB_CHUNKS = chunks(_IMAGE_LIST_CHUNK, _COLS)
            for j, sub_chunk in enumerate(_SUB_CHUNKS):
                for _ID, img in enumerate(sub_chunk):
                    try:
                        case_id = img
                        img_path = os.path.join(_ROOT_IMG, img)
                        if not os.path.exists(img_path):
                            img_path = os.path.join(_ROOT_IMG, img.lower())
                        if img_path is not None:
                            slide.shapes.add_picture(img_path, _lefts_vol[_ID],_tops[j], height=_ht)
                        txBox = slide.shapes.add_textbox(_lefts_txt[_ID],_tops_txt[j],Cm(1),Cm(0.25))
                        tf = txBox.text_frame
                        p = tf.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s\nActual: %s\nPred: %s\n' % (case_id, 
                                                            _CASE_LIST.loc[case_id, 'True'], 
                                                            _CASE_LIST.loc[case_id, 'Pred'])
                        font = run.font
                        font.name = 'Calibri'
                        font.size = Pt(8)
                        font.bold = True
                        font.color.theme_color = MSO_THEME_COLOR.ACCENT_5
                    except Exception as e:
                        print(e)
                        txBox = slide.shapes.add_textbox(_lefts_txt[_ID],_tops_txt[j],Cm(1),Cm(0.25))
                        tf = txBox.text_frame
                        p = tf.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s\nActual: %s\nPred: %s\n%s' % (case_id, 
                                                            _CASE_LIST.loc[case_id, 'True'], 
                                                            _CASE_LIST.loc[case_id, 'Pred'])
                        font = run.font
                        font.name = 'Calibri'
                        font.size = Pt(8)
                        font.bold = True
                        font.color.theme_color = MSO_THEME_COLOR.ACCENT_5
                        with open(_LOG_FILE, 'a') as f:
                            f.writelines('%s\n' % (img))
        prs.save(_ROOT_TGT_PPT+pptx_path)

if __name__ == '__main__':
    main()